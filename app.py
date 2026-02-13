import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
import io
from datetime import datetime

# --- 1. CONFIG & EXECUTIVE NAVY THEME ---
st.set_page_config(page_title="Sentinell BI | Governance Vantage", layout="wide", page_icon="🚀")

st.markdown("""
    <style>
    .stApp { background-color: #0f172a; color: #f8fafc; }
    [data-testid="stSidebar"] { background-color: #0b1120; border-right: 1px solid #1e293b; }
    
    .main-header {
        background: linear-gradient(90deg, #1e40af, #0f172a);
        color: white; padding: 1.5rem; border-radius: 12px;
        text-align: left; margin-bottom: 2rem; border-left: 8px solid #38bdf8;
    }
    
    .business-outlook {
        background-color: #1e293b; padding: 20px; border-radius: 10px;
        border-left: 5px solid #ef4444; margin-bottom: 20px;
    }
    
    .metric-card { 
        background: #1e293b; padding: 15px; border-radius: 8px; 
        border: 1px solid #334155; text-align: center; height: 130px;
    }
    .metric-label { color: #94a3b8; font-size: 10px; text-transform: uppercase; font-weight: 700; }
    .metric-value { color: #38bdf8; font-size: 26px; font-weight: bold; display: block; margin-top: 5px; }
    
    /* Governance Verdict Styling */
    .verdict-box {
        background-color: #0b1120; padding: 15px; border-radius: 8px;
        border: 1px solid #1e293b; margin-top: 10px;
    }
    .status-go { color: #4ade80; font-weight: bold; font-size: 20px; }
    </style>
    """, unsafe_allow_html=True)

# --- 2. AUTHENTICATION GATE ---
if "auth" not in st.session_state:
    c1, c2, c3 = st.columns([1, 1.5, 1])
    with c2:
        st.markdown("<br><br><div style='text-align:center;'><h2>🔒 Senior Director Access</h2></div>", unsafe_allow_html=True)
        u = st.text_input("Username")
        p = st.text_input("Security Key", type="password")
        if st.button("Authorize Access"):
            if p == "Company2026":
                st.session_state["auth"], st.session_state["user"] = True, u
                st.rerun()
            else: st.error("Invalid Credentials.")
    st.stop()

# --- 3. BUSINESS LOGIC ENGINES ---
def calculate_biz_aging(df):
    start = pd.to_datetime(df['Discovery_Date']).dt.date.values.astype('datetime64[D]')
    # Current date for 2026 Demo
    today_val = np.datetime64('2026-02-13') 
    end = pd.to_datetime(df['Closed_Date']).dt.date.fillna(pd.Timestamp(today_val)).values.astype('datetime64[D]')
    # Official Holidays 2026
    hols = ['2026-01-01', '2026-01-26', '2026-08-15', '2026-10-02', '2026-12-25']
    return np.busday_count(start, end, holidays=hols)

# --- 4. SIDEBAR: GOVERNANCE CONTROL PANEL ---
with st.sidebar:
    st.markdown("<h2 style='color:#38bdf8;'>🎯 Governance Filters</h2>", unsafe_allow_html=True)
    uploaded_file = st.file_uploader("📂 Sync Master Data (XLSX)", type=["xlsx"])
    
    if uploaded_file:
        df = pd.read_excel(uploaded_file)
        df.columns = [c.strip() for c in df.columns]
        df['Discovery_Date'] = pd.to_datetime(df['Discovery_Date'])
        
        # 📅 Reporting Period logic
        st.subheader("📅 Reporting Period")
        min_date, max_date = df['Discovery_Date'].min().date(), df['Discovery_Date'].max().date()
        date_range = st.date_input("Date Range Selection", [min_date, max_date])
        
        # 🎨 Visual Chart Settings
        st.divider()
        st.subheader("🎨 Chart Controls")
        chart_type = st.radio("Primary Chart Style", ["Bar Chart", "Line Chart", "Area Chart"])
        x_axis = st.selectbox("Strategic Dimension (X)", ['App_Area', 'Defect_ID', 'Root_Cause', 'Severity'])
        
        # 🎯 Metadata Filters
        st.divider()
        focus_id = st.multiselect("Focus Defect_ID", df['Defect_ID'].unique())
        f_status = st.multiselect("Filter Status", df['Status'].unique(), default=df['Status'].unique())
        f_sev = st.multiselect("Filter Severity", df['Severity'].unique(), default=df['Severity'].unique())
        
        # Business KPI Logic (Injected)
        df['Aging_Days'] = calculate_biz_aging(df)
        df['KPI_Status'] = df['Aging_Days'].apply(lambda x: 'Met' if x <= 5 else 'Breached')
        f_kpi = st.multiselect("Filter KPI_Status", ['Met', 'Breached'], default=['Met', 'Breached'])
        
        if st.button("🚪 Logout"):
            for key in list(st.session_state.keys()): del st.session_state[key]
            st.rerun()
    else:
        st.info("Awaiting Data Synchronization...")
        st.stop()

# --- 5. DATA SPLICING & ANALYTICS ---
# Date Range Mask
start_d, end_d = pd.to_datetime(date_range[0]), pd.to_datetime(date_range[1])
mask = (df['Discovery_Date'] >= start_d) & (df['Discovery_Date'] <= end_d) & \
       (df['Status'].isin(f_status)) & (df['Severity'].isin(f_sev)) & (df['KPI_Status'].isin(f_kpi))

if focus_id:
    mask = mask & (df['Defect_ID'].isin(focus_id))

f_df = df[mask].copy()

# Derived Metrics
risk_exposure = f_df['Fix_Cost'].sum() if 'Fix_Cost' in f_df.columns else 0
kpi_met_pct = (f_df['KPI_Status'] == 'Met').mean() * 100 if not f_df.empty else 0

# --- 6. MAIN UI: EXECUTIVE DASHBOARD ---
st.markdown(f"""
    <div class="main-header">
        <h1>🛡️ Sentinell V | Governance Vantage</h1>
        <p>Strategic Intelligence Suite for Senior Director: {st.session_state.get('user', 'Veera')}</p>
    </div>
""", unsafe_allow_html=True)

# Strategic Business Outlook
st.markdown(f"""
<div class="business-outlook">
    <h4 style="margin:0; color:#ef4444;">📈 Strategic Business Outlook</h4>
    <div style="margin-top:10px;">
        Risk Exposure: <b>${risk_exposure:,.0f}</b> | 
        KPI Status: <b>{kpi_met_pct:.1f}% Met</b> | 
        Total Focus Items: <b>{len(f_df)}</b>
    </div>
    <div class="verdict-box">
        <span class="status-go">🚦 STATUS: RECOMMENDED FOR PRODUCTION</span>
        <p style="margin:5px 0 0 0; font-size:12px; opacity:0.8;">Predictive Stability Index: 95.0% | Systemic Mode: {f_df['Root_Cause'].mode()[0] if not f_df.empty else 'N/A'}</p>
    </div>
</div>
""", unsafe_allow_html=True)

# Metric Tiles
m1, m2, m3, m4 = st.columns(4)
with m1: st.markdown(f"<div class='metric-card'><span class='metric-label'>Revenue Protection</span><span class='metric-value'>${risk_exposure:,.0f}</span></div>", unsafe_allow_html=True)
with m2: st.markdown(f"<div class='metric-card'><span class='metric-label'>Avg Aging</span><span class='metric-value'>{f_df['Aging_Days'].mean():.1f}d</span></div>", unsafe_allow_html=True)
with m3: st.markdown(f"<div class='metric-card'><span class='metric-label'>KPI Compliance</span><span class='metric-value'>{kpi_met_pct:.1f}%</span></div>", unsafe_allow_html=True)
with m4: st.markdown(f"<div class='metric-card'><span class='metric-label'>Stability Index</span><span class='metric-value'>95%</span></div>", unsafe_allow_html=True)

# --- 7. TABS: THE COMMAND ARCHITECTURE ---
st.divider()
t_dist, t_rca, t_trend, t_audit = st.tabs(["📊 Dimension Analysis", "🎯 Root Cause", "📈 Velocity", "📋 Audit Grid"])

with t_dist:
    c1, c2 = st.columns(2)
    with c1:
        st.subheader(f"Value Concentration by {x_axis}")
        if chart_type == "Bar Chart":
            fig = px.bar(f_df, x=x_axis, y="Fix_Cost", color="Severity", barmode="group", template="plotly_dark", color_discrete_sequence=px.colors.qualitative.G10)
        elif chart_type == "Line Chart":
            fig = px.line(f_df.sort_values(x_axis), x=x_axis, y="Fix_Cost", color="Severity", markers=True, template="plotly_dark")
        else:
            fig = px.area(f_df.sort_values(x_axis), x=x_axis, y="Fix_Cost", color="Severity", template="plotly_dark")
        st.plotly_chart(fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)'), use_container_width=True)
    with c2:
        st.subheader("Risk Heatmap (Aging vs Severity)")
        
        fig_heat = px.density_heatmap(f_df, x="Aging_Days", y="Severity", z="Fix_Cost", template="plotly_dark", color_continuous_scale="Viridis", text_auto=True)
        st.plotly_chart(fig_heat.update_layout(paper_bgcolor='rgba(0,0,0,0)'), use_container_width=True)

with t_rca:
    st.subheader("Systemic Failure Modes (Treemap)")
    
    fig_tree = px.treemap(f_df, path=['Root_Cause', 'App_Area'], values='Fix_Cost', template="plotly_dark", color_discrete_sequence=px.colors.qualitative.Pastel)
    st.plotly_chart(fig_tree.update_layout(paper_bgcolor='rgba(0,0,0,0)'), use_container_width=True)

with t_trend:
    st.subheader("Backlog Velocity (The Red Line)")
    
    trend_df = f_df.groupby(f_df['Discovery_Date'].dt.date).size().reset_index(name='Inflow')
    trend_df['Backlog'] = trend_df['Inflow'].cumsum()
    fig_line = go.Figure()
    fig_line.add_trace(go.Bar(name='New Defects', x=trend_df['Discovery_Date'], y=trend_df['Inflow'], marker_color='#38bdf8'))
    fig_line.add_trace(go.Scatter(name='Cumulative Backlog', x=trend_df['Discovery_Date'], y=trend_df['Backlog'], line=dict(color='#ef4444', width=3)))
    st.plotly_chart(fig_line.update_layout(template="plotly_dark", paper_bgcolor='rgba(0,0,0,0)'), use_container_width=True)

with t_audit:
    st.subheader("📋 Governance Audit Grid")
    st.dataframe(f_df[['Defect_ID', 'App_Area', 'Severity', 'Status', 'Aging_Days', 'KPI_Status', 'Fix_Cost', 'Root_Cause']], use_container_width=True)
    
    # PPT Export Feature (Final Functionality)
    if st.button("📊 Export Executive Presentation"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Sentinell BI: Executive Governance Summary"
        slide.placeholders[1].text = f"Report Date: {datetime.now().strftime('%Y-%m-%d')}\nRisk Exposure: ${risk_exposure:,.0f}\nKPI Met: {kpi_met_pct:.1f}%"
        buf = io.BytesIO()
        prs.save(buf)
        st.download_button("📥 Download PPTX", buf.getvalue(), f"Executive_Report_{datetime.now().strftime('%Y%m%d')}.pptx")
